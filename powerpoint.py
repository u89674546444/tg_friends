from pptx import Presentation
from pptx.util import Inches
import os
import tempfile

# Путь к шаблону PowerPoint
template_path = '/Users/nikolajusakov/PycharmProjects/PythonProject/Презентация.pptx'

# Путь к корневому каталогу
root_directory = '/Users/nikolajusakov/PycharmProjects/PythonProject/reports'

# Функция для создания презентации по шаблону
def create_presentation_from_folder(folder_path, template_path):
    # Открываем шаблон PowerPoint
    prs = Presentation(template_path)

    # Пути к изображениям и текстовому файлу
    image1_path = os.path.join(folder_path, 'до.jpg')
    image2_path = os.path.join(folder_path, 'после.jpg')
    text_file_path = os.path.join(folder_path, 'report.txt')

    # Проверка существования файлов
    if not os.path.exists(image1_path):
        print(f"Ошибка: файл {image1_path} не найден.")
        return None
    if not os.path.exists(image2_path):
        print(f"Ошибка: файл {image2_path} не найден.")
        return None
    if not os.path.exists(text_file_path):
        print(f"Ошибка: файл {text_file_path} не найден.")
        return None

    # Читаем текст из файла
    with open(text_file_path, 'r', encoding='utf-8') as file:
        text_content = file.read()

    # Проходим по всем слайдам в презентации
    for slide in prs.slides:
        # Проходим по всем фигурам на слайде
        for shape in slide.shapes:
            # Проверяем, является ли фигура текстовым полем
            if shape.has_text_frame:
                # Заменяем метку {Text} на текст из файла
                if "{Text}" in shape.text:
                    print("Найдена метка {Text}. Заменяем на текст из файла.")
                    shape.text = shape.text.replace("{Text}", text_content)

                # Заменяем метку {Image1} на изображение
                if "{Image1}" in shape.text:
                    print("Найдена метка {Image1}. Вставляем изображение.")
                    # Удаляем текстовое поле с меткой
                    text_frame = shape.text_frame
                    text_frame.clear()  # Очищаем текстовое поле
                    # Добавляем изображение на место метки
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    try:
                        slide.shapes.add_picture(image1_path, left, top, width, height)
                        print("Изображение 1 успешно вставлено.")
                    except Exception as e:
                        print(f"Ошибка при вставке изображения 1: {e}")

                # Заменяем метку {Image2} на изображение
                if "{Image2}" in shape.text:
                    print("Найдена метка {Image2}. Вставляем изображение.")
                    # Удаляем текстовое поле с меткой
                    text_frame = shape.text_frame
                    text_frame.clear()  # Очищаем текстовое поле
                    # Добавляем изображение на место метки
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    try:
                        slide.shapes.add_picture(image2_path, left, top, width, height)
                        print("Изображение 2 успешно вставлено.")
                    except Exception as e:
                        print(f"Ошибка при вставке изображения 2: {e}")

    # Сохраняем презентацию в текущей папке
    presentation_name = os.path.basename(folder_path) + '_presentation.pptx'
    output_path = os.path.join(folder_path, presentation_name)
    prs.save(output_path)
    print(f"Презентация сохранена как {output_path}")
    return output_path

# Функция для объединения презентаций
def merge_presentations(presentation_paths, output_path):
    merged_prs = Presentation()

    for presentation_path in presentation_paths:
        prs = Presentation(presentation_path)
        for slide in prs.slides:
            # Копируем слайды в объединённую презентацию
            slide_layout = merged_prs.slide_layouts[6]  # Пустой слайд
            new_slide = merged_prs.slides.add_slide(slide_layout)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Копируем текстовое поле
                    new_shape = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    new_shape.text = shape.text
                elif shape.shape_type == 13:  # Изображение
                    # Сохраняем изображение во временный файл
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp_file:
                        tmp_file.write(shape.image.blob)
                        tmp_file_path = tmp_file.name

                    # Вставляем изображение из временного файла
                    new_slide.shapes.add_picture(
                        tmp_file_path, shape.left, shape.top, shape.width, shape.height
                    )

                    # Удаляем временный файл
                    os.unlink(tmp_file_path)

    # Сохраняем объединённую презентацию
    merged_prs.save(output_path)
    print(f"Объединённая презентация сохранена как {output_path}")

# Рекурсивный обход папок
def process_directory(directory, template_path):
    presentation_paths = []

    for root, dirs, files in os.walk(directory):
        # Если это конечная папка (нет вложенных папок)
        if not dirs:
            print(f"Обработка папки: {root}")
            presentation_path = create_presentation_from_folder(root, template_path)
            if presentation_path:
                presentation_paths.append(presentation_path)

    # Объединяем все презентации в одну
    if presentation_paths:
        # Получаем путь к предпоследней папке
        parent_folder = os.path.dirname(os.path.dirname(presentation_paths[0]))
        output_path = os.path.join(parent_folder, 'merged_presentation.pptx')
        merge_presentations(presentation_paths, output_path)

        # Удаляем отдельные презентации (опционально)
        for presentation_path in presentation_paths:
            os.remove(presentation_path)
            print(f"Удалена презентация: {presentation_path}")

# Запуск обработки
process_directory(root_directory, template_path)